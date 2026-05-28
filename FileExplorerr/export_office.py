#!/usr/bin/env python3
# export_office.py — generador Office/PDF para FileExplorerr
# Uso: python export_office.py <formato> <csv_input> <output_path> <titulo>

import sys, csv, os, datetime

def progress(pct):
    """Reporta progreso al proceso padre (C# lo lee por stdout)."""
    print(f"PROGRESS:{pct}", flush=True)

def leer_csv(path):
    with open(path, 'r', encoding='utf-8-sig', newline='') as f:
        reader = csv.reader(f)
        all_rows = list(reader)
    return (all_rows[0] if all_rows else [],
            all_rows[1:] if len(all_rows) > 1 else [])

# ── EXCEL ─────────────────────────────────────────────────────────────────────
def exportar_xlsx(headers, rows, output, titulo):
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    from openpyxl.cell import WriteOnlyCell

    progress(25)
    wb = Workbook(write_only=True)
    ws = wb.create_sheet(titulo[:31] or "Datos")

    hdr_fill = PatternFill("solid", fgColor="1A3A4A")
    hdr_font = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    hdr_aln  = Alignment(horizontal="center", vertical="center", wrap_text=False)
    par_fill = PatternFill("solid", fgColor="E8F4F8")
    imp_fill = PatternFill("solid", fgColor="FFFFFF")
    thin     = Side(border_style="thin", color="BBCCDD")
    brd      = Border(left=thin, right=thin, top=thin, bottom=thin)
    dat_font = Font(name="Calibri", size=10)

    progress(30)
    hdr_cells = []
    for h in headers:
        c = WriteOnlyCell(ws, value=h)
        c.fill = hdr_fill; c.font = hdr_font
        c.alignment = hdr_aln; c.border = brd
        hdr_cells.append(c)
    ws.append(hdr_cells)

    total = len(rows)
    for ri, row in enumerate(rows):
        fill = par_fill if ri % 2 == 0 else imp_fill
        cells = []
        for val in row:
            c = WriteOnlyCell(ws, value=val)
            c.fill = fill; c.font = dat_font; c.border = brd
            cells.append(c)
        ws.append(cells)
        if ri % 3000 == 0 and total > 0:
            progress(30 + int(65 * ri / total))

    progress(95)
    wb.save(output)
    progress(100)

# ── WORD ──────────────────────────────────────────────────────────────────────
def exportar_docx(headers, rows, output, titulo):
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    # Limitar a 5000 filas — Word no es una hoja de cálculo
    MAX = 5000
    truncated = len(rows) > MAX
    rows = rows[:MAX]

    progress(20)
    doc = Document()
    section = doc.sections[0]
    landscape = len(headers) > 6
    if landscape:
        section.page_width   = Cm(29.7)
        section.page_height  = Cm(21.0)
        section.left_margin  = Cm(1.2)
        section.right_margin = Cm(1.2)
    else:
        section.page_width   = Cm(21.0)
        section.page_height  = Cm(29.7)
        section.left_margin  = Cm(2.0)
        section.right_margin = Cm(2.0)
    section.top_margin    = Cm(1.5)
    section.bottom_margin = Cm(1.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(titulo)
    run.font.size = Pt(20); run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3A, 0x4A)

    sub_txt = f"{len(rows):,} de {len(rows):,} registros  ·  {len(headers)} columnas"
    if truncated:
        sub_txt = f"Primeras {MAX:,} filas  ·  {len(headers)} columnas  (dataset completo: usa Excel o PDF)"
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(sub_txt)
    r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0x48, 0xCA, 0xB4)
    doc.add_paragraph()

    progress(30)
    ncols  = len(headers)
    avail  = 25.0 if landscape else 17.0
    col_w  = max(1.2, avail / ncols)
    table  = doc.add_table(rows=1, cols=ncols)
    table.style = 'Table Grid'

    def set_cell_bg(cell, hex_color):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement('w:shd')
        shd.set(qn('w:val'),   'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'),  hex_color)
        tcPr.append(shd)

    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        cell.width = Cm(col_w)
        cell.text  = str(h)
        run = cell.paragraphs[0].runs[0]
        run.font.bold = True; run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        set_cell_bg(cell, '1A3A4A')

    total = len(rows)
    for ri, row in enumerate(rows):
        bg  = 'E8F4F8' if ri % 2 == 0 else 'FFFFFF'
        tr  = table.add_row()
        for ci, val in enumerate(row[:ncols]):
            cell = tr.cells[ci]
            cell.width = Cm(col_w)
            cell.text  = str(val) if val else ''
            if cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].font.size = Pt(8)
            set_cell_bg(cell, bg)
        if ri % 500 == 0 and total > 0:
            progress(30 + int(65 * ri / total))

    progress(95)
    doc.save(output)
    progress(100)

# ── POWERPOINT ────────────────────────────────────────────────────────────────
def exportar_pptx(headers, rows, output, titulo):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    MAX_ROWS = 500
    MAX_COLS = 10
    truncated = len(rows) > MAX_ROWS
    rows    = rows[:MAX_ROWS]
    headers = headers[:MAX_COLS]
    for i, r in enumerate(rows):
        rows[i] = r[:MAX_COLS]

    progress(15)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_bg(slide, hex_color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color)

    def add_txbox(slide, text, l, t, w, h, size=18, bold=False,
                  color="FFFFFF", align="center"):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        r  = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(color)

    # Portada
    s0 = prs.slides.add_slide(blank)
    set_bg(s0, "1A3A4A")
    add_txbox(s0, titulo,
        Inches(0.5), Inches(2.2), Inches(12.33), Inches(1.6),
        size=34, bold=True, color="FFFFFF")
    total_r = len(rows)
    note = (f"Mostrando {total_r} de {total_r} filas · {len(headers)} columnas"
            if not truncated else
            f"Primeras {MAX_ROWS} filas mostradas · {len(headers)} columnas")
    add_txbox(s0, note,
        Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.8),
        size=15, color="48CAB4")
    add_txbox(s0, datetime.datetime.now().strftime("%d/%m/%Y %H:%M"),
        Inches(0.5), Inches(4.9), Inches(12.33), Inches(0.5),
        size=12, color="8899AA")

    progress(25)
    FPP = 18
    ncols = len(headers)
    total_slides = max(1, -(-len(rows) // FPP))

    for s in range(total_slides):
        ini = s * FPP
        fin = min(ini + FPP, len(rows))
        slide = prs.slides.add_slide(blank)
        set_bg(slide, "0F2030")

        add_txbox(slide,
            f"{titulo}  —  filas {ini+1}–{fin}  ({s+1}/{total_slides})",
            Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5),
            size=12, bold=True, color="48CAB4", align="left")

        chunk  = [headers] + rows[ini:fin]
        n_rows = len(chunk)

        left   = Inches(0.3); top    = Inches(0.72)
        width  = Inches(12.7); height = Inches(6.58)
        tbl    = slide.shapes.add_table(n_rows, ncols, left, top, width, height).table

        col_w = int(width / ncols)
        for ci in range(ncols): tbl.columns[ci].width = col_w

        for ri, row_data in enumerate(chunk):
            is_hdr = ri == 0
            bg = (RGBColor.from_string("1A3A4A") if is_hdr else
                  (RGBColor.from_string("1E3040") if ri % 2 == 1
                   else RGBColor.from_string("263545")))
            fg = RGBColor(255,255,255) if is_hdr else RGBColor(0xD8,0xE8,0xF0)
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)[:26] if val else ''
                para = cell.text_frame.paragraphs[0]
                run  = para.runs[0] if para.runs else para.add_run()
                run.font.size  = Pt(8) if is_hdr else Pt(7)
                run.font.bold  = is_hdr
                run.font.color.rgb = fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg

        progress(25 + int(70 * (s + 1) / total_slides))

    progress(95)
    prs.save(output)
    progress(100)

# ── PDF ───────────────────────────────────────────────────────────────────────
def exportar_pdf(headers, rows, output, titulo):
    from reportlab.lib.pagesizes import A4, landscape as rl_landscape
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                     Paragraph, Spacer, PageBreak)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER

    styles     = getSampleStyleSheet()
    title_sty  = ParagraphStyle('T', parent=styles['Normal'],
        fontSize=16, fontName='Helvetica-Bold',
        textColor=colors.HexColor('#1A3A4A'), alignment=TA_CENTER, spaceAfter=4)
    sub_sty    = ParagraphStyle('S', parent=styles['Normal'],
        fontSize=9, fontName='Helvetica',
        textColor=colors.HexColor('#48CAB4'), alignment=TA_CENTER, spaceAfter=10)

    progress(20)
    doc = SimpleDocTemplate(output, pagesize=rl_landscape(A4),
        leftMargin=1.5*cm, rightMargin=1.5*cm,
        topMargin=1.5*cm, bottomMargin=1.5*cm)

    story = []
    story.append(Paragraph(titulo, title_sty))
    story.append(Paragraph(
        f"{len(rows):,} registros  ·  {len(headers)} columnas  ·  "
        f"{datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}",
        sub_sty))

    # Anchos de columna (muestra de 300 filas)
    ncols   = len(headers)
    maxch   = [max(len(str(h)), 4) for h in headers]
    sample  = rows[::max(1, len(rows)//300)]
    for row in sample:
        for i, v in enumerate(row[:ncols]):
            maxch[i] = max(maxch[i], min(len(str(v)), 32))
    total_ch = sum(maxch) or 1
    avail_w  = 25.4 * cm
    col_widths = [max(0.8*cm, (ch/total_ch)*avail_w) for ch in maxch]

    CHUNK = 600
    total  = len(rows)
    chunks = max(1, -(-total // CHUNK))

    progress(30)
    for ch in range(chunks):
        ini  = ch * CHUNK
        fin  = min(ini + CHUNK, total)
        data = [[str(h)[:32] for h in headers]]
        for row in rows[ini:fin]:
            data.append([str(v)[:32] if v else '' for v in row[:ncols]])

        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        tbl.setStyle(TableStyle([
            ('BACKGROUND',    (0,0),(-1,0),  colors.HexColor('#1A3A4A')),
            ('TEXTCOLOR',     (0,0),(-1,0),  colors.white),
            ('FONTNAME',      (0,0),(-1,0),  'Helvetica-Bold'),
            ('FONTSIZE',      (0,0),(-1,0),  7),
            ('ALIGN',         (0,0),(-1,0),  'CENTER'),
            ('FONTNAME',      (0,1),(-1,-1), 'Helvetica'),
            ('FONTSIZE',      (0,1),(-1,-1), 6),
            ('TEXTCOLOR',     (0,1),(-1,-1), colors.HexColor('#1A1A2A')),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),
             [colors.HexColor('#E8F4F8'), colors.white]),
            ('GRID',          (0,0),(-1,-1), 0.3, colors.HexColor('#AABBCC')),
            ('LINEABOVE',     (0,0),(-1,0),  1,   colors.HexColor('#1A3A4A')),
            ('LINEBELOW',     (0,0),(-1,0),  1,   colors.HexColor('#1A3A4A')),
            ('TOPPADDING',    (0,0),(-1,-1), 2),
            ('BOTTOMPADDING', (0,0),(-1,-1), 2),
            ('LEFTPADDING',   (0,0),(-1,-1), 3),
            ('RIGHTPADDING',  (0,0),(-1,-1), 3),
        ]))
        story.append(tbl)
        if ch < chunks - 1:
            story.append(PageBreak())
        progress(30 + int(65 * (ch + 1) / chunks))

    progress(95)
    doc.build(story)
    progress(100)

# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    if len(sys.argv) < 5:
        print("Uso: export_office.py <xlsx|docx|pptx|pdf> <csv> <output> <titulo>")
        sys.exit(1)

    fmt, csv_path, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    titulo = sys.argv[4] if len(sys.argv) > 4 else "Exportación"

    progress(5)
    headers, rows = leer_csv(csv_path)
    progress(15)

    try:
        if   fmt == 'xlsx': exportar_xlsx(headers, rows, out_path, titulo)
        elif fmt == 'docx': exportar_docx(headers, rows, out_path, titulo)
        elif fmt == 'pptx': exportar_pptx(headers, rows, out_path, titulo)
        elif fmt == 'pdf':  exportar_pdf (headers, rows, out_path, titulo)
        else:
            print(f"Formato desconocido: {fmt}", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        import traceback
        traceback.print_exc()
        sys.exit(1)
